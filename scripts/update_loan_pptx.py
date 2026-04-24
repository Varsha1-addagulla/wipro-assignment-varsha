"""One-off script to update Copy of loan-risk-multiagent.pptx per review feedback.

Run: python scripts/update_loan_pptx.py
"""

from __future__ import annotations

import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.shapes.base import BaseShape
from pptx.util import Inches, Pt


def set_paragraphs(shape: BaseShape, lines: list[str]) -> None:
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0


def replace_first_text_paragraph(shape: BaseShape, new_text: str) -> None:
    if not shape.has_text_frame:
        return
    p = shape.text_frame.paragraphs[0]
    p.clear()
    p.text = new_text


def main() -> int:
    src = Path(
        r"C:\Users\varsh\Downloads\Copy of loan-risk-multiagent.pptx"
    )
    if not src.is_file():
        print("Expected file not found:", src, file=sys.stderr)
        return 1

    bak = src.with_suffix(".pptx.bak")
    shutil.copy2(src, bak)
    print("Backup:", bak)

    prs = Presentation(str(src))

    def slide0(n: int):
        return prs.slides[n - 1]

    s1 = slide0(1)
    set_paragraphs(
        s1.shapes[3],
        [
            "Nine core specialists + Intake, Planner, and Negotiator in a "
            "LangGraph pipeline. Deterministic gates around LLM output.",
            "Agency-style (Fannie Mae–inspired) threshold bands. Deployed on "
            "GCP Cloud Run.",
            "Agenda: problem and outcome  →  architecture and agents  →  control "
            "flow and security  →  deploy and walkthrough  →  next steps and Q&A",
        ],
    )

    s2 = slide0(2)
    set_paragraphs(
        s2.shapes[2],
        [
            "The problem: first-pass lending needs synchronized decisions across "
            "credit, income, employment, fraud, and debt. Agency-style (Fannie "
            "Mae–inspired) threshold bands show how those dimensions bind together, "
            "but a single all-in-one LLM call blurs the trade-offs and leaves no "
            "reliable audit path.",
        ],
    )
    out_lines = [
        "Why this matters: a faster, clearer first-pass call (approve, review, or "
        "reject) with a trace, so borrowers are not stuck waiting with no line of "
        "sight, and operations spend less time on obvious triage.",
        "Accepts 6 fields: name, loan amount, annual income, credit score, "
        "employment years, existing debt.",
        "Pipeline: Intake (enrich)  →  Planner (fast-reject or full path)  →  "
        "Consistency  →  five parallel LLM specialists  →  Debt  →  Critic  →  "
        "Negotiator when not a full approve  →  Report.",
        "Returns: APPROVED, HUMAN_REVIEW, or AUTO_REJECTED.",
        "Each /assess run is written to api_log: request id, latency, decision, "
        "confidence, and structured metadata for audit and debugging.",
    ]
    set_paragraphs(s2.shapes[6], out_lines)

    s3 = slide0(3)
    for sh in s3.shapes:
        if not sh.has_text_frame:
            continue
        t = (sh.text_frame.text or "").strip()
        if t == "All traffic terminates at Cloud Run (HTTPS, autoscaled, secret-injected).":
            sh.text_frame.paragraphs[0].text = (
                "Up-front: Intake and Planner. After the Critic: Negotiator, then "
                "Report. All traffic: HTTPS to Cloud Run (autoscaled, secret-injected)."
            )
            break

    s4 = slide0(4)
    replace_first_text_paragraph(
        s4.shapes[1], "The agent team: nine specialists and orchestration"
    )
    replace_first_text_paragraph(
        s4.shapes[3],
        "Two deterministic layers (plus the Planner) wrap the LLM work; the graph "
        "owns order and parallel execution. Gates and the critic can override model "
        "output where rules require it, never the reverse. Also in the same graph, "
        "but not in the table: Intake, Planner, and Negotiator — see the note below.",
    )
    s4.shapes.add_textbox(
        Inches(0.45),
        Inches(4.5),
        Inches(9.0),
        Inches(0.95),
    )
    box = s4.shapes[-1]
    btf = box.text_frame
    btf.clear()
    p0 = btf.paragraphs[0]
    p0.text = (
        "Also in the deployed graph: Intake (context before specialists), Planner "
        "(fast-reject path vs full LLM run), and Negotiator (counter-offer after the "
        "Critic). The table lists the five parallel specialists, Debt, Critic, and "
        "Report Writer, plus the Consistency and Fraud gates as shown."
    )
    p0.font.size = Pt(12)

    s5 = slide0(5)
    replace_first_text_paragraph(
        s5.shapes[5], "Intake, Planner, validate"
    )
    replace_first_text_paragraph(
        s5.shapes[6], "Intake, Planner, Consistency"
    )
    replace_first_text_paragraph(
        s5.shapes[7],
        "Intake enriches; Planner chooses a fast-reject or full run. Consistency is "
        "pure Python, runs early, and can hard-stop the graph before parallel LLM work.",
    )
    replace_first_text_paragraph(
        s5.shapes[23], "Decide, negotiate, report"
    )
    replace_first_text_paragraph(
        s5.shapes[24], "Critic, Negotiator, Report Writer"
    )
    replace_first_text_paragraph(
        s5.shapes[25],
        "Critic applies thresholds. Negotiator proposes a structured counter-offer "
        "when the vote is not a clean approval. Report Writer turns the run into a "
        "narrative suitable for review.",
    )
    s5.shapes[26].text_frame.paragraphs[0].text = (
        "State moves through a shared dict; agents do not call each other. LangGraph "
        "owns ordering, the parallel specialist stage, and when to short-circuit."
    )

    s6 = slide0(6)
    for sh in s6.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            t = (para.text or "")
            if "Fannie" in t and "rationalizes" in t:
                para.text = (
                    "The LLM receives pre-labeled credit and DTI bands from code. It does "
                    "not invent the cutoffs. It explains against the same Fannie-style "
                    "illustration used in the assignment. JSON-only, with regex repair "
                    "if the model adds fences."
                )
                break

    s8 = slide0(8)
    for sh in s8.shapes:
        if not sh.has_text_frame:
            continue
        t = (sh.text_frame.text or "")
        if "Only applicant name" in t or t.startswith("Data handling"):
            replace_first_text_paragraph(
                sh,
                "Data handling: api_log rows for audit, including applicant name. "
                "Structured logs redact api_key, authorization, and cookies.",
            )
            break

    s9 = slide0(9)
    for sh in s9.shapes:
        if not sh.has_text_frame:
            continue
        t = (sh.text_frame.text or "")
        if "SQLite" in t and "Cloud SQL" in t:
            p1 = sh.text_frame.paragraphs[1]
            p1.text = "SQLite by default. Cloud SQL in production when enabled."
            break

    s12 = slide0(12)
    replace_first_text_paragraph(
        s12.shapes[1], "Trade-offs, next steps, and Q&A"
    )
    set_paragraphs(
        s12.shapes[9],
        [
            "Golden-set evaluation — pinned cases and expected decision bands on "
            "every change.",
            "Tighter JSON output — use Groq JSON mode or a schema-validated path "
            "so regex fallbacks are rare.",
            "Observability — track p95 latency for full path vs planner fast-reject path.",
        ],
    )
    thanks = s12.shapes.add_textbox(
        Inches(0.5),
        Inches(5.0),
        Inches(9.0),
        Inches(0.6),
    )
    tp = thanks.text_frame.paragraphs[0]
    tp.text = "Thank you — questions?"
    tp.font.size = Pt(20)
    tp.font.bold = True
    tp.font.color.rgb = RGBColor(0x1A, 0x3A, 0x5C)

    prs.save(str(src))
    print("Updated:", src)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
